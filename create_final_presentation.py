from pathlib import Path
import zipfile
import xml.etree.ElementTree as ET

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


ROOT = Path(__file__).resolve().parent
OUT = ROOT / "SturzDetektion_Abschlusspraesentation_final.pptx"
ASSET_DIR = ROOT / "fall_realdata_project" / "outputs"
GEN_DIR = ROOT / "praesentation_assets"
GEN_DIR.mkdir(exist_ok=True)

NAVY = RGBColor(17, 24, 39)
MUTED = RGBColor(102, 112, 133)
BLUE = RGBColor(37, 99, 235)
TEAL = RGBColor(13, 148, 136)
GREEN = RGBColor(22, 163, 74)
ORANGE = RGBColor(245, 158, 11)
RED = RGBColor(220, 38, 38)
PURPLE = RGBColor(124, 58, 237)
BG = RGBColor(248, 250, 252)
WHITE = RGBColor(255, 255, 255)
LINE = RGBColor(226, 232, 240)
PALE_BLUE = RGBColor(239, 246, 255)
PALE_GREEN = RGBColor(240, 253, 244)
PALE_RED = RGBColor(254, 242, 242)
PALE_ORANGE = RGBColor(255, 247, 237)
PALE_TEAL = RGBColor(240, 253, 250)
PALE_PURPLE = RGBColor(245, 243, 255)

HEX = ["#2563eb", "#0d9488", "#f59e0b", "#dc2626", "#16a34a", "#7c3aed"]


def save_bar(path, labels, values, title, colors):
    fig, ax = plt.subplots(figsize=(8.6, 4.4), dpi=160)
    bars = ax.bar(labels, values, color=colors, width=0.58)
    ax.set_title(title, fontweight="bold", pad=14)
    ax.grid(axis="y", alpha=0.22)
    ax.spines[["top", "right"]].set_visible(False)
    ax.set_axisbelow(True)
    ax.set_ylim(0, max(values) * 1.18)
    for bar, value in zip(bars, values):
        ax.text(
            bar.get_x() + bar.get_width() / 2,
            value + max(values) * 0.025,
            f"{value:,.0f}".replace(",", "."),
            ha="center",
            fontweight="bold",
        )
    fig.tight_layout()
    fig.savefig(path, facecolor="white")
    plt.close(fig)


def save_pie(path):
    labels = ["Ruhiger Alltag", "Normale Bewegung", "Sturzähnlich OK", "Sturz"]
    values = [5749, 10000, 10000, 10000]
    fig, ax = plt.subplots(figsize=(7.2, 4.4), dpi=160)
    wedges, _, autotexts = ax.pie(
        values,
        autopct=lambda p: f"{p:.1f}%",
        startangle=90,
        colors=[HEX[4], HEX[0], HEX[2], HEX[3]],
        pctdistance=0.75,
        wedgeprops=dict(width=0.42, edgecolor="white"),
    )
    ax.legend(wedges, labels, loc="center left", bbox_to_anchor=(0.80, 0.5), frameon=False)
    ax.set_title("Klassenverteilung nach Balancing", fontweight="bold", pad=14)
    for text in autotexts:
        text.set_fontweight("bold")
    fig.tight_layout()
    fig.savefig(path, facecolor="white")
    plt.close(fig)


def save_metrics(path):
    labels = ["Ruhig", "Normal", "Sturzähnl.", "Sturz"]
    precision = [77.7, 97.1, 96.9, 99.3]
    recall = [99.3, 91.1, 92.4, 95.7]
    f1 = [87.2, 94.0, 94.6, 97.5]
    x = np.arange(len(labels))
    width = 0.24
    fig, ax = plt.subplots(figsize=(8.9, 4.5), dpi=160)
    ax.bar(x - width, precision, width, label="Precision", color=HEX[0])
    ax.bar(x, recall, width, label="Recall", color=HEX[2])
    ax.bar(x + width, f1, width, label="F1", color=HEX[4])
    ax.set_ylim(65, 103)
    ax.set_xticks(x, labels)
    ax.set_ylabel("%")
    ax.set_title("Klassenauswertung im Testsplit", fontweight="bold", pad=14)
    ax.grid(axis="y", alpha=0.22)
    ax.spines[["top", "right"]].set_visible(False)
    ax.legend(frameon=False, ncol=3, loc="upper center", bbox_to_anchor=(0.5, -0.08))
    fig.tight_layout()
    fig.savefig(path, facecolor="white")
    plt.close(fig)


def save_confusion(path):
    matrix = np.array(
        [[679, 0, 5, 0], [69, 1071, 33, 3], [83, 30, 1437, 5], [43, 2, 8, 1194]]
    )
    labels = ["Ruhig", "Normal", "Sturzähnl.", "Sturz"]
    fig, ax = plt.subplots(figsize=(6.6, 5.2), dpi=160)
    image = ax.imshow(matrix, cmap="Blues")
    ax.set_title("Konfusionsmatrix", fontweight="bold", pad=14)
    ax.set_xlabel("Vorhergesagte Klasse")
    ax.set_ylabel("Echte Klasse")
    ax.set_xticks(np.arange(4), labels, rotation=25, ha="right")
    ax.set_yticks(np.arange(4), labels)
    max_value = matrix.max()
    for row in range(4):
        for col in range(4):
            ax.text(
                col,
                row,
                str(matrix[row, col]),
                ha="center",
                va="center",
                color="white" if matrix[row, col] > max_value * 0.45 else "#111827",
                fontweight="bold",
            )
    fig.colorbar(image, ax=ax, fraction=0.046, pad=0.04)
    fig.tight_layout()
    fig.savefig(path, facecolor="white")
    plt.close(fig)


def save_thresholds(path):
    labels = ["Sturz-Wahrscheinlichkeit", "Impact", "Nachbewegung", "Bestätigte Fenster"]
    display = ["> 90 %", "> 28 m/s²", "< 1,8", "≥ 2"]
    widths = [90, 70, 45, 55]
    colors = [HEX[3], HEX[2], HEX[1], HEX[0]]
    fig, ax = plt.subplots(figsize=(8.5, 4.2), dpi=160)
    bars = ax.barh(labels, widths, color=colors, height=0.5)
    ax.set_xlim(0, 100)
    ax.set_xticks([])
    ax.spines[:].set_visible(False)
    ax.set_title("Schwellwerte der Sicherheitslogik", fontweight="bold", pad=14)
    for bar, text in zip(bars, display):
        ax.text(bar.get_width() + 2, bar.get_y() + bar.get_height() / 2, text, va="center", fontweight="bold")
    fig.tight_layout()
    fig.savefig(path, facecolor="white")
    plt.close(fig)


def build_charts():
    charts = {
        "sources": GEN_DIR / "datenquellen_final.png",
        "classes": GEN_DIR / "klassen_final.png",
        "split": GEN_DIR / "split_final.png",
        "metrics": GEN_DIR / "metrics_final.png",
        "confusion": GEN_DIR / "confusion_final.png",
        "thresholds": GEN_DIR / "thresholds_final.png",
    }
    save_bar(charts["sources"], ["SisFall", "UniMiB", "UP-Fall"], [19639, 8036, 8074], "Fenster pro Datensatz", HEX[:3])
    save_pie(charts["classes"])
    save_bar(charts["split"], ["Train", "Validation", "Test"], [25435, 5652, 4662], "Train / Validation / Test", [HEX[0], HEX[2], HEX[4]])
    save_metrics(charts["metrics"])
    save_confusion(charts["confusion"])
    save_thresholds(charts["thresholds"])
    return charts


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def text(slide, content, x, y, w, h, size=18, bold=False, color=NAVY, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.04)
    frame.margin_right = Inches(0.04)
    frame.margin_top = Inches(0.02)
    frame.margin_bottom = Inches(0.02)
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = content
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def bullets(slide, items, x, y, w, h, size=16):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.05)
    frame.margin_right = Inches(0.05)
    frame.margin_top = Inches(0.02)
    frame.margin_bottom = Inches(0.02)
    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = item
        paragraph.font.name = "Aptos"
        paragraph.font.size = Pt(size)
        paragraph.font.color.rgb = NAVY
        paragraph.space_after = Pt(7)
    return box


def card(slide, x, y, w, h, fill=WHITE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = LINE
    shape.line.width = Pt(0.75)
    return shape


def title(slide, heading, eyebrow=None, accent=BLUE):
    accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.13), Inches(7.5))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = accent
    accent_bar.line.fill.background()
    if eyebrow:
        text(slide, eyebrow.upper(), 0.72, 0.28, 5.5, 0.24, 9.5, True, accent)
    text(slide, heading, 0.70, 0.55, 11.8, 0.55, 26, True, NAVY)


def metric(slide, value, label, x, y, w, accent=BLUE):
    card(slide, x, y, w, 0.86)
    text(slide, value, x + 0.08, y + 0.10, w - 0.16, 0.30, 18, True, accent, PP_ALIGN.CENTER)
    text(slide, label, x + 0.08, y + 0.50, w - 0.16, 0.22, 9.5, False, MUTED, PP_ALIGN.CENTER)


def picture(slide, path, x, y, w, h):
    with Image.open(path) as image:
        image_w, image_h = image.size
    box_ratio = w / h
    image_ratio = image_w / image_h
    if image_ratio > box_ratio:
        new_w = w
        new_h = w / image_ratio
        new_x = x
        new_y = y + (h - new_h) / 2
    else:
        new_h = h
        new_w = h * image_ratio
        new_x = x + (w - new_w) / 2
        new_y = y
    return slide.shapes.add_picture(str(path), Inches(new_x), Inches(new_y), width=Inches(new_w), height=Inches(new_h))


def add_footer():
    for index, slide in enumerate(prs.slides, 1):
        text(slide, str(index), 12.55, 7.02, 0.35, 0.18, 8.5, False, MUTED, PP_ALIGN.RIGHT)


def process_slide(heading, eyebrow, steps, body, accent=TEAL):
    slide = prs.slides.add_slide(BLANK)
    background(slide)
    title(slide, heading, eyebrow, accent)
    for index, (head, sub) in enumerate(steps):
        x = 0.75 + index * 2.45
        y = 1.55
        card(slide, x, y, 2.02, 1.15)
        text(slide, head, x + 0.12, y + 0.15, 1.78, 0.26, 12.5, True, accent, PP_ALIGN.CENTER)
        text(slide, sub, x + 0.13, y + 0.52, 1.76, 0.42, 9.5, False, MUTED, PP_ALIGN.CENTER)
        if index < len(steps) - 1:
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x + 2.05), Inches(y + 0.43), Inches(0.28), Inches(0.24))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = accent
            arrow.line.fill.background()
    card(slide, 0.85, 3.25, 11.6, 2.65)
    bullets(slide, body, 1.15, 3.55, 10.95, 1.95, 17)


def build_presentation():
    charts = build_charts()

    slide = prs.slides.add_slide(BLANK)
    background(slide)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.20))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLUE
    bar.line.fill.background()
    card(slide, 0.75, 0.95, 11.85, 5.5)
    text(slide, "ABSCHLUSSPRÄSENTATION", 1.15, 1.35, 4.4, 0.25, 10, True, BLUE)
    text(slide, "SturzDetektion", 1.08, 1.75, 8.6, 0.8, 42, True)
    text(slide, "KI-basierte Sturzerkennung mit Android-App und 1D-CNN", 1.15, 2.65, 9.8, 0.45, 19, False, MUTED)
    for i, (value, label) in enumerate([("Android", "Java-App"), ("150 x 6", "Inputfenster"), ("1D-CNN", "Modell"), ("93,97 %", "Test Accuracy")]):
        metric(slide, value, label, 1.15 + i * 2.75, 4.35, 2.25, BLUE)

    process_slide(
        "Projekt in einem Bild",
        "Überblick",
        [("Sensoren", "Accel + Gyro"), ("Fenster", "3 Sekunden"), ("KI-Modell", "TFLite 1D-CNN"), ("Regeln", "Sicherheitscheck"), ("Alarm", "Countdown")],
        [
            "Die App analysiert Bewegungsdaten lokal auf dem Smartphone.",
            "Das neuronale Netz bewertet jedes Zeitfenster mit vier Klassenwahrscheinlichkeiten.",
            "Eine zusätzliche Sicherheitslogik entscheidet, ob wirklich ein Alarm ausgelöst wird.",
        ],
        TEAL,
    )

    slide = prs.slides.add_slide(BLANK)
    background(slide)
    title(slide, "Problem und Ziel", "Motivation", RED)
    card(slide, 0.85, 1.35, 5.65, 4.7, PALE_RED)
    text(slide, "Problem", 1.15, 1.70, 2.2, 0.35, 20, True, RED)
    bullets(slide, ["Stürze können unbemerkt bleiben.", "Hilfe kommt oft verzögert.", "Alltagsbewegungen können Fehlalarme verursachen."], 1.15, 2.25, 4.95, 2.8, 18)
    card(slide, 6.85, 1.35, 5.65, 4.7, PALE_GREEN)
    text(slide, "Ziel", 7.15, 1.70, 2.2, 0.35, 20, True, GREEN)
    bullets(slide, ["Bewegungen kontinuierlich erfassen.", "Stürze lokal mit KI erkennen.", "Fehlalarme durch Sicherheitsregeln reduzieren.", "Im Ernstfall einen Countdown starten."], 7.15, 2.25, 4.95, 2.8, 18)

    slide = prs.slides.add_slide(BLANK)
    background(slide)
    title(slide, "App-Aufbau", "Android", BLUE)
    for x, heading, items, fill in [
        (0.85, "UI-Schicht", ["Dashboard", "Testmodus", "Notfallkarte"], PALE_BLUE),
        (4.65, "Laufzeitlogik", ["SensorEventListener", "Ringbuffer", "processDetection"], PALE_TEAL),
        (8.45, "KI-Schicht", ["TFLite Interpreter", "metadata.json", "4 Klassen"], PALE_PURPLE),
    ]:
        card(slide, x, 1.45, 3.55, 4.75, fill)
        text(slide, heading, x + 0.25, 1.75, 3.05, 0.32, 18, True, NAVY, PP_ALIGN.CENTER)
        bullets(slide, items, x + 0.45, 2.35, 2.65, 2.6, 17)
    text(slide, "Technik: Java, AppCompat, Material Components, TensorFlow Lite 2.16.1, minSdk 26", 0.95, 6.55, 11.8, 0.28, 12, False, MUTED, PP_ALIGN.CENTER)

    slide = prs.slides.add_slide(BLANK)
    background(slide)
    title(slide, "Funktionen der App", "Bedienung", GREEN)
    features = [
        ("Start/Stop", "Live-Überwachung aktivieren"),
        ("Dashboard", "Sensorquelle und Modellstatus"),
        ("Wahrscheinlichkeiten", "alle vier Klassen sichtbar"),
        ("Testmodus", "Szenarien simulieren"),
        ("Notfall", "30-Sekunden-Countdown"),
        ("Log", "Erkennung nachvollziehen"),
    ]
    for i, (head, body) in enumerate(features):
        x = 0.85 + (i % 3) * 4.05
        y = 1.5 + (i // 3) * 2.05
        card(slide, x, y, 3.55, 1.45)
        text(slide, head, x + 0.18, y + 0.22, 3.18, 0.26, 15, True, GREEN)
        text(slide, body, x + 0.18, y + 0.68, 3.18, 0.38, 12.5, False, MUTED)

    process_slide(
        "Sensorpipeline",
        "Live-Erkennung",
        [("Accel", "ax, ay, az"), ("Gyro", "gx, gy, gz"), ("Ringbuffer", "150 Samples"), ("Normierung", "Mean / Std"), ("Inferenz", "alle 400 ms")],
        [
            "Fenstergröße: 150 Samples = 3,0 Sekunden bei 50 Hz.",
            "Die App nutzt dieselbe Normalisierung wie das Trainingsscript.",
            "Beschleunigungs- und Gyro-Magnituden werden zusätzlich für Sicherheitsregeln berechnet.",
        ],
        TEAL,
    )

    slide = prs.slides.add_slide(BLANK)
    background(slide)
    title(slide, "Datengrundlage", "Trainingsdaten", BLUE)
    metric(slide, "35.749", "Fenster final", 0.85, 1.28, 2.35, BLUE)
    metric(slide, "3", "Datensätze", 3.35, 1.28, 2.0, BLUE)
    metric(slide, "50 Hz", "Zielrate", 5.50, 1.28, 2.0, BLUE)
    bullets(slide, ["SisFall, UniMiB und UP-Fall wurden vereinheitlicht.", "Alle Sequenzen werden auf 150 x 6 Fenster gebracht.", "Balancing verhindert eine zu starke Dominanz einzelner Klassen."], 0.95, 2.55, 4.6, 2.8, 16)
    picture(slide, charts["sources"], 5.9, 1.55, 6.35, 4.7)

    slide = prs.slides.add_slide(BLANK)
    background(slide)
    title(slide, "Klassenverteilung", "Nach Balancing", PURPLE)
    picture(slide, charts["classes"], 0.9, 1.35, 6.0, 4.85)
    card(slide, 7.15, 1.45, 5.05, 4.55)
    text(slide, "Vier Zielklassen", 7.5, 1.78, 3.6, 0.35, 19, True, PURPLE)
    bullets(slide, ["0 Ruhiger Alltag", "1 Normale Bewegung", "2 Sturzähnlich, aber OK", "3 Sturz"], 7.55, 2.35, 4.2, 2.3, 18)
    text(slide, "Klasse 2 ist wichtig, weil sie harte, aber ungefährliche Bewegungen abfängt.", 7.55, 5.15, 4.15, 0.55, 13, False, MUTED)

    slide = prs.slides.add_slide(BLANK)
    background(slide)
    title(slide, "Train / Validation / Test", "Subject Split", ORANGE)
    picture(slide, charts["split"], 0.9, 1.45, 6.4, 4.75)
    card(slide, 7.65, 1.55, 4.6, 4.5, PALE_ORANGE)
    text(slide, "Warum nach Personen splitten?", 7.95, 1.9, 4.0, 0.32, 18, True, ORANGE)
    bullets(slide, ["Keine identischen Personen in Train und Test.", "Geringeres Risiko für Data Leakage.", "Realistischere Aussage über Generalisierung."], 7.95, 2.45, 3.75, 2.2, 17)
    metric(slide, "60 / 15 / 10", "Subjects in Train / Val / Test", 8.15, 5.15, 3.55, ORANGE)

    slide = prs.slides.add_slide(BLANK)
    background(slide)
    title(slide, "Neuronales Netz", "1D-CNN Architektur", RED)
    steps = [("Input", "150 x 6"), ("Conv1D 32", "BN + ReLU"), ("Conv1D 64", "Pooling"), ("Conv1D 128", "Dropout"), ("Conv1D 128", "Features"), ("Avg + Max", "Pooling"), ("Dense 96", "Dropout"), ("Softmax", "4 Klassen")]
    for index, (head, body) in enumerate(steps):
        x = 0.65 + (index % 4) * 3.1
        y = 1.45 + (index // 4) * 2.0
        card(slide, x, y, 2.55, 1.25)
        text(slide, head, x + 0.12, y + 0.18, 2.3, 0.25, 13, True, RED, PP_ALIGN.CENTER)
        text(slide, body, x + 0.12, y + 0.58, 2.3, 0.28, 10.5, False, MUTED, PP_ALIGN.CENTER)
        if index not in [3, 7]:
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x + 2.58), Inches(y + 0.47), Inches(0.35), Inches(0.25))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = RED
            arrow.line.fill.background()
    metric(slide, "112.516", "Parameter", 2.2, 5.85, 2.35, RED)
    metric(slide, "111.620", "trainierbar", 5.1, 5.85, 2.35, RED)
    metric(slide, "4", "Ausgabeklassen", 8.0, 5.85, 2.35, RED)

    slide = prs.slides.add_slide(BLANK)
    background(slide)
    title(slide, "Training", "Setup", BLUE)
    for i, (value, label) in enumerate([("TensorFlow 2.21", "Framework"), ("Batch 128", "Batch Size"), ("60", "Epochen"), ("1e-3", "Learning Rate")]):
        metric(slide, value, label, 0.85 + i * 3.0, 1.35, 2.45, BLUE)
    card(slide, 0.9, 2.65, 5.45, 3.25, PALE_BLUE)
    text(slide, "Augmentation", 1.2, 2.95, 3.6, 0.35, 19, True, BLUE)
    bullets(slide, ["leichtes Rauschen", "Skalierung 0,92 bis 1,08", "Gyro-Dropout 30 %", "Time Shift ±10 Samples"], 1.2, 3.55, 4.4, 1.8, 16)
    card(slide, 6.75, 2.65, 5.45, 3.25, PALE_GREEN)
    text(slide, "Schutz vor Overfitting", 7.05, 2.95, 4.3, 0.35, 19, True, GREEN)
    bullets(slide, ["subject-stratified split", "Class Weights", "Dropout im Netz", "Validation während Training"], 7.05, 3.55, 4.4, 1.8, 16)

    slide = prs.slides.add_slide(BLANK)
    background(slide)
    title(slide, "Trainingsverlauf", "Accuracy und Loss", BLUE)
    picture(slide, ASSET_DIR / "real_training_accuracy.png", 0.8, 1.35, 5.85, 4.55)
    picture(slide, ASSET_DIR / "real_training_loss.png", 6.85, 1.35, 5.85, 4.55)
    text(slide, "Die Kurven dienen als Plausibilitätscheck: Accuracy steigt, Loss sinkt und stabilisiert sich.", 1.0, 6.35, 11.3, 0.35, 13, False, MUTED, PP_ALIGN.CENTER)

    slide = prs.slides.add_slide(BLANK)
    background(slide)
    title(slide, "Testergebnisse", "Gesamtleistung", GREEN)
    for i, (value, label, color) in enumerate([("93,97 %", "Accuracy", GREEN), ("0,2292", "Test Loss", BLUE), ("99,3 %", "Sturz Precision", RED), ("95,7 %", "Sturz Recall", ORANGE)]):
        metric(slide, value, label, 0.85 + i * 3.0, 1.4, 2.45, color)
    card(slide, 1.15, 3.0, 10.95, 2.65)
    bullets(slide, ["Das Modell erkennt die Klasse Sturz sehr zuverlässig.", "Die schwierigsten Fehler liegen zwischen Sturz und sturzähnlicher Bewegung.", "Genau deshalb nutzt die App zusätzlich Beschleunigung, Nachbewegung und Mehrfachfenster-Bestätigung."], 1.55, 3.45, 10.1, 1.55, 18)

    slide = prs.slides.add_slide(BLANK)
    background(slide)
    title(slide, "Klassenauswertung", "Precision / Recall / F1", GREEN)
    picture(slide, charts["metrics"], 0.8, 1.35, 8.15, 4.95)
    card(slide, 9.25, 1.55, 3.1, 4.35, PALE_GREEN)
    text(slide, "Interpretation", 9.55, 1.9, 2.5, 0.32, 17, True, GREEN)
    bullets(slide, ["Sturz: höchste Precision.", "Ruhig: sehr hoher Recall.", "Klasse 2 trennt harte Bewegungen von echten Stürzen."], 9.55, 2.45, 2.45, 2.15, 14.5)

    slide = prs.slides.add_slide(BLANK)
    background(slide)
    title(slide, "Konfusionsmatrix", "Testsplit", ORANGE)
    picture(slide, charts["confusion"], 0.95, 1.25, 6.45, 5.35)
    card(slide, 7.75, 1.55, 4.55, 4.65, PALE_ORANGE)
    text(slide, "Wichtigste Zahlen", 8.05, 1.9, 3.7, 0.32, 18, True, ORANGE)
    bullets(slide, ["1.194 von 1.247 Stürzen korrekt.", "Nur 5 sturzähnliche Bewegungen als Sturz vorhergesagt.", "Hauptproblem: Grenzfälle zwischen Sturz und sturzähnlich OK.", "Sicherheitsregeln sollen diese Fälle auffangen."], 8.05, 2.45, 3.75, 2.75, 15)

    slide = prs.slides.add_slide(BLANK)
    background(slide)
    title(slide, "Sicherheitslogik", "KI plus Regeln", RED)
    picture(slide, charts["thresholds"], 0.85, 1.45, 6.75, 4.4)
    card(slide, 7.95, 1.55, 4.45, 4.45, PALE_RED)
    text(slide, "Alarm nur bei Kombination", 8.25, 1.9, 3.85, 0.32, 17, True, RED)
    bullets(slide, ["Klasse 3 mit hoher Wahrscheinlichkeit", "starker Impact", "wenig Bewegung danach", "mindestens zwei bestätigte Fenster"], 8.25, 2.45, 3.65, 2.2, 15.5)
    text(slide, "Das verhindert, dass einzelne Peaks sofort einen Notfall auslösen.", 8.25, 5.25, 3.55, 0.42, 12.5, False, MUTED)

    slide = prs.slides.add_slide(BLANK)
    background(slide)
    title(slide, "Testmodus", "Simulationen", TEAL)
    scenarios = [("Echter Sturz", "Alarm"), ("Rempler", "kein Sturz"), ("Gehen", "Normal"), ("Joggen", "Normal"), ("Handy fällt", "Verdächtig"), ("Stolpern", "Grenzfall"), ("Treppe", "Normal"), ("Sturz low/high rot", "Alarmtest")]
    for i, (head, body) in enumerate(scenarios):
        x = 0.85 + (i % 4) * 3.05
        y = 1.45 + (i // 4) * 1.85
        card(slide, x, y, 2.55, 1.25, PALE_TEAL)
        text(slide, head, x + 0.12, y + 0.18, 2.3, 0.28, 12.5, True, TEAL, PP_ALIGN.CENTER)
        text(slide, body, x + 0.12, y + 0.65, 2.3, 0.25, 10.5, False, MUTED, PP_ALIGN.CENTER)
    card(slide, 1.05, 5.25, 11.25, 0.8)
    text(slide, "Der Testmodus zeigt erwartetes Ergebnis, KI-Sturzwahrscheinlichkeit, Sicherheitsregel und Gesamtbewertung.", 1.25, 5.50, 10.85, 0.25, 14, False, MUTED, PP_ALIGN.CENTER)

    process_slide(
        "Deployment",
        "Vom Training in die App",
        [("Keras", ".keras"), ("Export", ".tflite"), ("Assets", "Modell + Metadata"), ("Android", "Interpreter"), ("UI", "Ergebnis + Alarm")],
        [
            "Das App-Modell liegt als fall_detection_real_150x6_float32.tflite in app/src/main/assets.",
            "metadata.json speichert Mean und Standardabweichung für die identische Normalisierung.",
            "Die Inferenz läuft lokal auf dem Gerät und benötigt keine Cloud-Verbindung.",
        ],
        PURPLE,
    )

    slide = prs.slides.add_slide(BLANK)
    background(slide)
    title(slide, "Grenzen und Ausblick", "Weiterentwicklung", ORANGE)
    card(slide, 0.9, 1.35, 5.45, 4.9, PALE_ORANGE)
    text(slide, "Aktuelle Grenzen", 1.2, 1.7, 4.6, 0.35, 19, True, ORANGE)
    bullets(slide, ["Notruf ist simuliert.", "Keine dauerhafte Hintergrundüberwachung.", "Noch keine echte Nutzerstudie.", "Datensätze kommen aus verschiedenen Labor-Setups."], 1.2, 2.3, 4.4, 2.8, 16)
    card(slide, 6.85, 1.35, 5.45, 4.9, PALE_BLUE)
    text(slide, "Nächste Schritte", 7.15, 1.7, 4.6, 0.35, 19, True, BLUE)
    bullets(slide, ["mehr eigene Smartphone-Daten", "Background Service", "Kontakt-/Notrufintegration", "Vergleich mit LSTM/GRU/Transformer"], 7.15, 2.3, 4.4, 2.8, 16)

    slide = prs.slides.add_slide(BLANK)
    background(slide)
    title(slide, "Fazit", "Kernaussage", BLUE)
    card(slide, 1.05, 1.55, 11.25, 4.65)
    bullets(slide, ["Das Projekt ist durchgängig umgesetzt: Datenaufbereitung, Training, TFLite-Export und Android-App.", "Das 1D-CNN erreicht eine starke Testleistung und ist klein genug für lokale Smartphone-Inferenz.", "Die App kombiniert KI-Ausgabe mit Sicherheitsregeln, um Fehlalarme zu reduzieren.", "Damit ist SturzDetektion eine lauffähige Grundlage für KI-basierte Sturzerkennung im Alltag."], 1.55, 2.0, 10.25, 3.45, 20)

    add_footer()
    prs.save(OUT)


def validate():
    Presentation(str(OUT))
    with zipfile.ZipFile(OUT) as archive:
        bad = archive.testzip()
        if bad:
            raise RuntimeError(f"Defekter ZIP-Eintrag: {bad}")
        for name in archive.namelist():
            if name.endswith(".xml") or name.endswith(".rels"):
                ET.fromstring(archive.read(name))


if __name__ == "__main__":
    build_presentation()
    validate()
    print(OUT)
    print(f"Slides: {len(prs.slides)}")
    print(f"Size: {OUT.stat().st_size / 1024 / 1024:.2f} MB")
